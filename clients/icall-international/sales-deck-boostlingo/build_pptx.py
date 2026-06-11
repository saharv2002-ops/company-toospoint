"""Build the iCall × Boostlingo vendor pitch as a 16:9 PowerPoint backed by
the rendered slide images, with hyperlink hot-spots that mirror the HTML deck.
"""
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SRC_DIR = os.path.dirname(os.path.abspath(__file__))
PNG_DIR = os.path.join(SRC_DIR, 'slides_png')
OUT     = os.path.join(SRC_DIR, 'iCall_Boostlingo_Vendor_Deck.pptx')

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# (left, top, width, height, url) — fractions of slide
HOTSPOTS = {
    1: [
        # "Become A Vendor Partner" CTA pill (bottom-right)
        (0.665, 0.832, 0.270, 0.080, 'https://icallinternational.com'),
    ],
    7: [
        # tel / mailto / web contact rows (left column)
        (0.07, 0.69, 0.30, 0.06, 'tel:+18182394744'),
        (0.07, 0.76, 0.30, 0.06, 'mailto:info@icallinternational.com'),
        (0.07, 0.83, 0.30, 0.06, 'https://icallinternational.com'),
        # "Schedule Vendor Review" CTA button (centre)
        (0.42, 0.78, 0.30, 0.110, 'https://icallinternational.com'),
        # logo (right)
        (0.79, 0.78, 0.15, 0.110, 'https://icallinternational.com'),
    ],
}


def _xml(s):
    wrapped = f'<root xmlns:a="{A_NS}" xmlns:p="{P_NS}">{s}</root>'
    return list(etree.fromstring(wrapped))[0]


def make_clean_hotspot(slide, left, top, width, height, url, tooltip=''):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sp = shp._element

    style_el = sp.find(qn('p:style'))
    if style_el is not None:
        sp.remove(style_el)

    txBody = sp.find(qn('p:txBody'))
    if txBody is not None:
        sp.remove(txBody)
    sp.append(_xml(
        '<p:txBody>'
        '  <a:bodyPr wrap="none" rtlCol="0" anchor="ctr"/>'
        '  <a:lstStyle/>'
        '  <a:p><a:endParaRPr lang="en-US"/></a:p>'
        '</p:txBody>'
    ))

    spPr = sp.find(qn('p:spPr'))
    for fill_tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        for el in spPr.findall(qn(fill_tag)):
            spPr.remove(el)
    for ln in spPr.findall(qn('a:ln')):
        spPr.remove(ln)
    spPr.append(_xml(
        '<a:solidFill>'
        '  <a:srgbClr val="FFFFFF"><a:alpha val="0"/></a:srgbClr>'
        '</a:solidFill>'
    ))
    spPr.append(_xml('<a:ln><a:noFill/></a:ln>'))

    shp.click_action.hyperlink.address = url

    nvSpPr = sp.find(qn('p:nvSpPr'))
    cNvPr  = nvSpPr.find(qn('p:cNvPr'))
    cNvPr.set('descr', tooltip or url)
    hlink = cNvPr.find(qn('a:hlinkClick'))
    if hlink is not None and tooltip:
        hlink.set('tooltip', tooltip)

    return shp


def add_slide_with_image(prs, png_path, hotspots=None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    slide.shapes.add_picture(png_path, 0, 0,
                             width=prs.slide_width, height=prs.slide_height)

    if hotspots:
        for left_f, top_f, w_f, h_f, url in hotspots:
            make_clean_hotspot(
                slide,
                int(prs.slide_width  * left_f),
                int(prs.slide_height * top_f),
                int(prs.slide_width  * w_f),
                int(prs.slide_height * h_f),
                url,
                tooltip=url,
            )


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    pngs = sorted(p for p in os.listdir(PNG_DIR) if p.endswith('.png'))
    for i, name in enumerate(pngs, start=1):
        add_slide_with_image(prs, os.path.join(PNG_DIR, name), HOTSPOTS.get(i))

    cp = prs.core_properties
    cp.title    = 'iCall International — Vendor Pitch for Boostlingo'
    cp.author   = 'iCall International'
    cp.subject  = 'LSP vendor partnership — compliance, QA, languages'
    cp.keywords = 'interpretation, OPI, VRI, Boostlingo, vendor, compliance, SDN, OIG, ESS, Alta, LanguageStat, QA'

    prs.save(OUT)
    print('Wrote:', OUT)


if __name__ == '__main__':
    main()
