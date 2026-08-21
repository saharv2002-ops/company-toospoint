"""Build the Simplified Group x Remoto discovery deck as a 16:9 PowerPoint
backed by the rendered slide images."""
import os
from pptx import Presentation
from pptx.util import Inches

SRC_DIR = os.path.dirname(os.path.abspath(__file__))
PNG_DIR = os.path.join(SRC_DIR, 'slides_png')
OUT     = os.path.join(SRC_DIR, 'Deck.pptx')

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_slide_with_image(prs, png_path):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(png_path, 0, 0, width=prs.slide_width, height=prs.slide_height)


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    pngs = sorted(p for p in os.listdir(PNG_DIR) if p.endswith('.png'))
    for name in pngs:
        add_slide_with_image(prs, os.path.join(PNG_DIR, name))

    cp = prs.core_properties
    cp.title    = 'Simplified Group x Remoto — Partnership Discovery'
    cp.author   = 'Simplified Group BPO'
    cp.subject  = 'Discovery deck for Remoto partnership call'
    cp.keywords = 'Remoto, agency partner, interpretation, OPI, VRI, platform'

    prs.save(OUT)
    print('Wrote:', OUT)


if __name__ == '__main__':
    main()
