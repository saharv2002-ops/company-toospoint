"""Build the SGP sales deck as a 16:9 PowerPoint backed by the rendered slide
images, with hyperlink hot-spots that mirror the HTML deck's clickable areas.

Produces clean OOXML compatible with modern PowerPoint (2016+), Office 365,
Keynote, and Google Slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SRC_DIR = os.path.dirname(os.path.abspath(__file__))
PNG_DIR = os.path.join(SRC_DIR, 'slides_png')
OUT     = os.path.join(SRC_DIR, 'SGP_Sales_Deck.pptx')

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# (left, top, width, height, url) — fractions of slide
HOTSPOTS = {
    1: [
        (0.683, 0.832, 0.260, 0.083, 'https://simplifiedbpo.com'),
    ],
    4: [
        (0.0625, 0.512, 0.215, 0.382, 'https://simplifiedbpo.com/opi/'),
        (0.297,  0.512, 0.215, 0.382, 'https://simplifiedbpo.com/vri/'),
        (0.531,  0.512, 0.215, 0.382, 'https://simplifiedbpo.com/'),
        (0.766,  0.512, 0.215, 0.382, 'https://simplifiedbpo.com/customer-service/'),
    ],
    11: [
        (0.07, 0.69, 0.27, 0.075, 'tel:+13025500551'),
        (0.07, 0.78, 0.27, 0.075, 'mailto:recruitmentdept@simplifiedbpo.com'),
        (0.07, 0.87, 0.27, 0.075, 'https://simplifiedbpo.com'),
        (0.39, 0.78, 0.30, 0.110, 'https://simplifiedbpo.com'),
        (0.74, 0.78, 0.21, 0.110, 'https://simplifiedbpo.com'),
    ],
}


def _xml(s):
    """Parse XML fragment with the full namespace declarations."""
    wrapped = (
        f'<root xmlns:a="{A_NS}" xmlns:p="{P_NS}">{s}</root>'
    )
    return list(etree.fromstring(wrapped))[0]


def make_clean_hotspot(slide, left, top, width, height, url, tooltip=''):
    """Add an invisible-but-clickable rectangle backed by clean OOXML.

    - Transparent solid fill (alpha 0) so the hit-area is honoured by
      every PowerPoint version, but renders nothing on screen.
    - No line, no text, no preset shape style.
    - Single shape-level click action; no run-level hyperlink.
    """
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sp = shp._element

    # Drop python-pptx's auto-added preset style (the accent1 reference)
    style_el = sp.find(qn('p:style'))
    if style_el is not None:
        sp.remove(style_el)

    # Drop the default text body and replace with an empty one
    txBody = sp.find(qn('p:txBody'))
    if txBody is not None:
        sp.remove(txBody)
    sp.append(_xml(
        '<p:txBody>'
        '  <a:bodyPr wrap="none" rtlCol="0" anchor="ctr"/>'
        '  <a:lstStyle/>'
        '  <a:p><a:endParaRPr lang="en-US"/></a:p>'
        '</p:txBody>'
    ))

    # Replace fill with transparent (alpha=0) — required so PowerPoint
    # registers the click area on shapes that have no visible body.
    spPr = sp.find(qn('p:spPr'))
    for fill_tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        for el in spPr.findall(qn(fill_tag)):
            spPr.remove(el)
    for ln in spPr.findall(qn('a:ln')):
        spPr.remove(ln)
    spPr.append(_xml(
        '<a:solidFill>'
        '  <a:srgbClr val="FFFFFF"><a:alpha val="0"/></a:srgbClr>'
        '</a:solidFill>'
    ))
    spPr.append(_xml('<a:ln><a:noFill/></a:ln>'))

    # Single, clean shape-level hyperlink (click_action manages rels properly)
    shp.click_action.hyperlink.address = url

    # Tooltip + accessibility name
    nvSpPr = sp.find(qn('p:nvSpPr'))
    cNvPr  = nvSpPr.find(qn('p:cNvPr'))
    cNvPr.set('descr', tooltip or url)
    # The hlinkClick already exists under cNvPr from click_action.
    hlink = cNvPr.find(qn('a:hlinkClick'))
    if hlink is not None and tooltip:
        hlink.set('tooltip', tooltip)

    return shp


def add_slide_with_image(prs, png_path, hotspots=None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Full-bleed picture (no hyperlink on the picture itself)
    pic = slide.shapes.add_picture(png_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Hot-spot overlays
    if hotspots:
        for left_f, top_f, w_f, h_f, url in hotspots:
            make_clean_hotspot(
                slide,
                int(prs.slide_width  * left_f),
                int(prs.slide_height * top_f),
                int(prs.slide_width  * w_f),
                int(prs.slide_height * h_f),
                url,
                tooltip=url,
            )


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    pngs = sorted(p for p in os.listdir(PNG_DIR) if p.endswith('.png'))
    for i, name in enumerate(pngs, start=1):
        add_slide_with_image(prs, os.path.join(PNG_DIR, name), HOTSPOTS.get(i))

    cp = prs.core_properties
    cp.title    = 'Simplified Group BPO — Sales Deck'
    cp.author   = 'Simplified Group BPO'
    cp.subject  = 'Capabilities & Service Overview'
    cp.keywords = 'OPI, VRI, Sales, Customer Service, BPO, interpretation'

    prs.save(OUT)
    print('Wrote:', OUT)


if __name__ == '__main__':
    main()
