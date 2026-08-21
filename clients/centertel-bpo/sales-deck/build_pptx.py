"""Assemble Center Tel BPO sales deck PPTX from rendered PNGs.

Designed for cross-platform compatibility: PowerPoint on Windows and macOS,
Keynote on macOS/iPadOS/iOS, and Google Slides. Each slide is a single
full-bleed picture; a per-slide primary hyperlink is attached directly to the
picture (portable across renderers) instead of using invisible shape overlays,
which Mac PowerPoint and Keynote often render as opaque rectangles.
"""
import os
from pptx import Presentation
from pptx.util import Inches

SRC_DIR = os.path.dirname(os.path.abspath(__file__))
PNG_DIR = os.path.join(SRC_DIR, 'slides_png')
OUT     = os.path.join(SRC_DIR, 'CenterTel_BPO_Sales_Deck.pptx')

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# One primary hyperlink per slide (attached to the whole slide picture).
SLIDE_LINKS = {
    1:  'https://center-tel.com',
    4:  'https://center-tel.com',
    12: 'https://center-tel.com',
}


def add_slide_with_image(prs, png_path, url=None):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    pic = slide.shapes.add_picture(png_path, 0, 0,
                                   width=prs.slide_width, height=prs.slide_height)
    if url:
        pic.click_action.hyperlink.address = url


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    pngs = sorted(p for p in os.listdir(PNG_DIR) if p.endswith('.png'))
    for i, name in enumerate(pngs, start=1):
        add_slide_with_image(prs, os.path.join(PNG_DIR, name), SLIDE_LINKS.get(i))
    cp = prs.core_properties
    cp.title    = 'Center Tel BPO — Sales Deck'
    cp.author   = 'Center Tel BPO'
    cp.subject  = 'Capabilities & Service Overview'
    cp.keywords = 'BPO, customer experience, call center, interpretation, OPI, VRI, sales support, outbound, back office, security monitoring, multilingual'
    prs.save(OUT)
    print('Wrote:', OUT)


if __name__ == '__main__':
    main()
